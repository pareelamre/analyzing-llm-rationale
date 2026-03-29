"""
Build: Context Structuring for LLM Reasoning
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Palette ───────────────────────────────────────────────────────────────────
C_BG    = RGBColor(0x0D, 0x11, 0x17)
C_PANEL = RGBColor(0x13, 0x1C, 0x2B)
C_BLUE  = RGBColor(0x00, 0x7A, 0xFF)
C_TEAL  = RGBColor(0x00, 0xC2, 0xA8)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY = RGBColor(0xC8, 0xD4, 0xE4)
C_MGRAY = RGBColor(0x78, 0x8C, 0xA8)
C_DGRAY = RGBColor(0x2A, 0x38, 0x50)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, line=None, lw=0):
    s = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h,
        size=16, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    r.font.name      = "Calibri"
    return tb


def bullets(slide, items, l, t, w, h, size=15, color=C_LGRAY, gap=6):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(gap)
        r = p.add_run()
        r.text           = f"  •  {item}"
        r.font.size      = Pt(size)
        r.font.color.rgb = color
        r.font.name      = "Calibri"


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.1, C_PANEL)
    txt(slide, title, 0.5, 0.18, 12, 0.6,
        size=26, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, 0.5, 0.72, 12, 0.35,
            size=13, color=C_TEAL)


def accent_bar(slide):
    rect(slide, 0, 1.1, 13.33, 0.04, C_BLUE)


def slide_num(slide, n):
    txt(slide, str(n), 12.8, 7.1, 0.4, 0.3,
        size=10, color=C_MGRAY, align=PP_ALIGN.RIGHT)


# ── Slide 1 — Title ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 2.5, 13.33, 2.8, C_PANEL)
rect(s, 0, 2.5, 0.06, 2.8, C_BLUE)

txt(s, "Structuring Context for Better LLM Reasoning",
    0.7, 2.8, 11.5, 1.2,
    size=36, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
txt(s, "How the way you frame information changes what an LLM can do with it",
    0.7, 4.15, 11.5, 0.6,
    size=18, italic=True, color=C_TEAL, align=PP_ALIGN.LEFT)
slide_num(s, 1)


# ── Slide 2 — The Core Idea ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "The Core Idea",
       "Not how much context — but which kind, and how it's shaped")
accent_bar(s)

bullets(s, [
    "Giving an LLM a raw question leaves it with nothing to anchor on",
    "Dumping too much unstructured information overwhelms it",
    "The real lever is deciding what goes in, and in what form",
], 0.6, 1.4, 12, 2.0, size=18, gap=14)

txt(s,
    "This project tests that idea directly — by building a dataset where each question\n"
    "is paired with carefully structured context, then measuring how different structures\n"
    "affect LLM accuracy and calibration.",
    0.6, 3.7, 12, 1.8, size=16, color=C_LGRAY)
slide_num(s, 2)


# ── Slide 3 — Experimental Pipeline ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Experimental Pipeline")
accent_bar(s)

# Pipeline boxes
BOX_W = 5.5
BOX_H = 0.9
CX    = (13.33 - BOX_W) / 2   # horizontally centered

labels = [
    ("Metaculus  —  1,580 questions",
     "binary · resolved · resolution criteria + news articles",
     C_BLUE),
    ("8 Prompt Variants",
     "each targets a different reasoning component",
     C_TEAL),
    ("5 LLMs  —  proprietary + open-weight",
     "output: { answer · confidence · rationale }",
     C_TEAL),
    ("Evaluation",
     "Accuracy  ·  Brier Score  ·  ECE  ·  Human rationale quality",
     RGBColor(0x00, 0xD4, 0x7E)),
]

tops = [1.4, 2.7, 4.0, 5.3]

for (title, sub, col), top in zip(labels, tops):
    rect(s, CX, top, BOX_W, BOX_H, C_PANEL, line=col, lw=1.2)
    txt(s, title, CX + 0.2, top + 0.08, BOX_W - 0.3, 0.45,
        size=14, bold=True, color=col)
    txt(s, sub,   CX + 0.2, top + 0.48, BOX_W - 0.3, 0.38,
        size=11, color=C_MGRAY)

# Arrow connectors
arrow_x = CX + BOX_W / 2 - 0.015
for top in tops[:-1]:
    rect(s, arrow_x, top + BOX_H, 0.03, 0.3, C_DGRAY)

slide_num(s, 3)


# ── Slide 4 — How Context Is Structured ──────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "How Context Is Structured")
accent_bar(s)

txt(s,
    "Each question is augmented with layers of context — question, background,\n"
    "news evidence, and resolution criteria.",
    0.6, 1.4, 12, 1.0, size=17, color=C_LGRAY)

# Layer stack — left to right progression
layers = ["Question", "+ Background", "+ News Evidence", "+ Resolution Criteria"]
colors = [C_DGRAY, C_PANEL,
          RGBColor(0x1A, 0x30, 0x50), RGBColor(0x00, 0x50, 0xAA)]
lw = 2.6
gap = 0.18
start_l = (13.33 - (len(layers) * lw + (len(layers) - 1) * gap)) / 2
top_l = 2.7

for i, (label, col) in enumerate(zip(layers, colors)):
    lx = start_l + i * (lw + gap)
    rect(s, lx, top_l, lw, 1.1, col, line=C_BLUE, lw=0.8)
    txt(s, label, lx + 0.1, top_l + 0.3, lw - 0.2, 0.55,
        size=14, bold=(i == 0), color=C_WHITE, align=PP_ALIGN.CENTER)

# Arrow between each
for i in range(len(layers) - 1):
    ax = start_l + (i + 1) * (lw + gap) - gap / 2 - 0.02
    txt(s, "▶", ax, top_l + 0.35, 0.25, 0.4,
        size=13, color=C_BLUE, align=PP_ALIGN.CENTER)

txt(s,
    "Key design choice:  curated and concise over raw and complete.\n"
    "Each layer adds a specific type of signal, not just more text.",
    0.6, 4.3, 12, 1.0, size=16, italic=True, color=C_TEAL)
slide_num(s, 4)


# ── Slide 5 — The 8 Prompt Variants ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "The 8 Prompt Variants",
       "Testing how the reasoning is framed — not just what context is included")
accent_bar(s)

txt(s,
    "Each variant asks the model to structure its rationale around a different component —\n"
    "the predicted outcome, key conditions, reasoning type, uncertainty, temporal deadline, and more.",
    0.6, 1.4, 12, 1.1, size=16, color=C_LGRAY)

txt(s,
    "The goal is to isolate which dimensions of structured reasoning\n"
    "actually improve forecasting accuracy and calibration.",
    0.6, 5.5, 12, 0.9, size=15, italic=True, color=C_TEAL)

# Variant chips
variants = [
    "V0  Neutral Baseline",
    "V1  Predicted Event",
    "V2  Key Attribute",
    "V3  Reasoning Type",
    "V5  Key Conditions",
    "V6  Step-by-Step",
    "V7  Uncertainty Language",
    "V8  Temporal Anchors",
]
chip_w, chip_h = 3.0, 0.55
cols_n = 4
gap_x, gap_y = 0.22, 0.18
total_w = cols_n * chip_w + (cols_n - 1) * gap_x
start_x = (13.33 - total_w) / 2
start_y = 2.7

for i, v in enumerate(variants):
    col_i = i % cols_n
    row_i = i // cols_n
    cx = start_x + col_i * (chip_w + gap_x)
    cy = start_y + row_i * (chip_h + gap_y)
    rect(s, cx, cy, chip_w, chip_h, C_PANEL, line=C_TEAL, lw=0.8)
    txt(s, v, cx + 0.15, cy + 0.1, chip_w - 0.2, chip_h - 0.1,
        size=13, color=C_WHITE)

slide_num(s, 5)


# ── Slide 6 — What the Results Show ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "What the Results Show")
accent_bar(s)

findings = [
    "Some structures improve accuracy; others improve how humans perceive the reasoning — these are not the same thing",
    "Temporal grounding is the most consistent and transferable improvement across all models",
    "Forcing explicit uncertainty language can hurt calibration rather than help it",
    "The right structure depends on the model and what you are optimizing for",
]
bullets(s, findings, 0.6, 1.4, 12, 4.0, size=17, gap=18)

slide_num(s, 6)


# ── Slide 7 — Why This Generalizes ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Why This Generalizes")
accent_bar(s)

txt(s,
    "The same structuring choices appear in any domain where an LLM must reason under uncertainty.",
    0.6, 1.4, 12, 0.7, size=17, color=C_LGRAY)

domains = [
    ("Medicine",  "Key conditions for diagnosis\nDeadline for action"),
    ("Law",       "Conditions that must hold\nType of reasoning being applied"),
    ("Finance",   "Credibility of the signal\nTemporal horizon"),
]
box_w = 3.6
gap_d = 0.5
start_d = (13.33 - len(domains) * box_w - (len(domains) - 1) * gap_d) / 2

for i, (domain, desc) in enumerate(domains):
    dx = start_d + i * (box_w + gap_d)
    rect(s, dx, 2.4, box_w, 2.6, C_PANEL, line=C_BLUE, lw=1.0)
    txt(s, domain, dx + 0.2, 2.55, box_w - 0.3, 0.55,
        size=18, bold=True, color=C_BLUE)
    txt(s, desc, dx + 0.2, 3.15, box_w - 0.3, 1.6,
        size=14, color=C_LGRAY)

txt(s,
    "The components tested here are not forecasting-specific.\n"
    "They are the building blocks of structured reasoning in any high-stakes context.",
    0.6, 5.4, 12, 0.9, size=15, italic=True, color=C_TEAL)

slide_num(s, 7)


# ── Slide 8 — Takeaway ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "The Takeaway")
accent_bar(s)

takeaways = [
    "Context structure is a design decision, not an afterthought",
    "Accuracy, calibration, and interpretability do not always point to the same structure",
    "This research gives a principled basis for deciding what to tell an LLM — not just which LLM to use",
]
bullets(s, takeaways, 0.6, 1.5, 12, 3.5, size=20, gap=24)

slide_num(s, 8)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "context_structuring.pptx"
prs.save(out)
print(f"Saved → {out}")
